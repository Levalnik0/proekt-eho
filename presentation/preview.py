"""
Проверка собранной презентации: читаем готовый .pptx и рисуем его в HTML
по тем же координатам. Так видно переполнение текста и выходы за край —
LibreOffice в системе нет, а глазами по коду такие вещи не поймать.
"""
import base64
import html
import io
import os
import sys

from pptx import Presentation
from pptx.util import Emu

args = [a for a in sys.argv[1:] if not a.startswith('--')]
SRC = args[0] if args else 'Проект_Эхо_презентация.pptx'
OUT = 'preview.html'

# --only 7 оставляет в превью один слайд. Нужно, чтобы снять его отдельным
# кадром: снимок делается с целой страницы, и на общем листе слайд мелкий.
ONLY = next(
    (int(a.split('=')[1]) for a in sys.argv[1:] if a.startswith('--only=')),
    None,
)
# --scale=2 рисует слайд вдвое крупнее. Съехавший на пару пикселей маркер
# в обычном размере просто не разглядеть.
SCALE = next(
    (float(a.split('=')[1]) for a in sys.argv[1:] if a.startswith('--scale=')),
    1.0,
)
PX = 96 * SCALE  # 1 дюйм = 96 px

# Cambria и Calibri — шрифты Microsoft, в macOS их нет. Подставляем то, что
# есть в системе и не уже оригинала: если текст влез здесь, в PowerPoint он
# влезет тем более. Точность ширины при этом теряется — на глаз проверяем
# композицию, а не последний пиксель строки.
FONTS = {
    'Cambria': 'Cambria, Caladea, Georgia, serif',
    'Calibri': 'Calibri, Carlito, Arial, sans-serif',
}

prs = Presentation(SRC)
SW = prs.slide_width / 914400 * PX
SH = prs.slide_height / 914400 * PX

problems = []
parts = []

def emu(v):
    return (v or 0) / 914400 * PX

for idx, slide in enumerate(prs.slides, 1):
    if ONLY and idx != ONLY:
        continue
    # Фон берём из самого слайда, иначе тёмные слайды в превью выглядят
    # светлыми и белый текст на них не виден. Фон бывает и картинкой.
    style = 'background:#F2F4F6'
    try:
        bg = slide._element.find(
            './/{http://schemas.openxmlformats.org/presentationml/2006/main}bg'
        )
        if bg is not None:
            clr = bg.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
            blip = bg.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}blip')
            if blip is not None:
                rid = blip.get(
                    '{http://schemas.openxmlformats.org/officeDocument/2006/'
                    'relationships}embed'
                )
                data = slide.part.related_part(rid).blob
                b64 = base64.b64encode(data).decode()
                style = (
                    f'background-image:url(data:image/png;base64,{b64});'
                    'background-size:cover'
                )
            elif clr is not None:
                style = f'background:#{clr.get("val")}'
    except Exception:
        pass
    items = []
    for sh in slide.shapes:
        x, y = emu(sh.left), emu(sh.top)
        w, h = emu(sh.width), emu(sh.height)

        if x < -1 or y < -1 or x + w > SW + 1 or y + h > SH + 1:
            problems.append(f'слайд {idx}: «{sh.shape_type}» выходит за край слайда')

        if sh.shape_type == 13:  # picture
            blob = sh.image.blob
            b64 = base64.b64encode(blob).decode()
            items.append(
                f'<img src="data:image/png;base64,{b64}" style="left:{x}px;top:{y}px;'
                f'width:{w}px;height:{h}px">'
            )
            continue

        fill = ''
        try:
            if sh.fill.type is not None and sh.fill.type == 1:
                fill = f'background:#{sh.fill.fore_color.rgb}'
        except Exception:
            pass
        prst = ''
        try:
            g = sh._element.find(
                './/{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom'
            )
            if g is not None:
                prst = g.get('prst') or ''
        except Exception:
            pass

        if fill:
            radius = ''
            if prst == 'roundRect':
                radius = f'border-radius:{0.14 * PX:.0f}px;'
            elif prst == 'ellipse':
                radius = 'border-radius:50%;'
            items.append(
                f'<div class="shape" style="left:{x}px;top:{y}px;width:{w}px;'
                f'height:{h}px;{fill};{radius}"></div>'
            )

        if sh.has_text_frame and sh.text_frame.text.strip():
            # Коды python-pptx: 1 — верх, 3 — середина, 4 — низ
            va = {3: 'center', 4: 'flex-end'}.get(
                int(sh.text_frame.vertical_anchor) if sh.text_frame.vertical_anchor else 0,
                'flex-start',
            )
            runs = []
            size, color, bold, italic, align = 14, '232A32', False, False, 'left'
            family = FONTS['Calibri']
            spacing = None
            for p in sh.text_frame.paragraphs:
                if p.alignment is not None:
                    # Коды python-pptx: 1 — влево, 2 — по центру, 3 — вправо
                    align = {2: 'center', 3: 'right'}.get(int(p.alignment), 'left')
                if p.line_spacing is not None:
                    spacing = p.line_spacing
                for r in p.runs:
                    if r.font.size:
                        size = r.font.size.pt
                    if r.font.name:
                        family = FONTS.get(r.font.name, r.font.name)
                    if r.font.bold:
                        bold = True
                    if r.font.italic:
                        italic = True
                    try:
                        if r.font.color and r.font.color.rgb:
                            color = str(r.font.color.rgb)
                    except Exception:
                        pass
                    runs.append(html.escape(r.text))
            text = ' '.join(runs)
            # Размер шрифта в файле задан в пунктах, а рисуем в пикселях.
            # 13 пунктов — это 17 пикселей: без пересчёта превью врало на
            # четверть, и по нему нельзя было судить о выравнивании.
            size_px = size / 72 * PX
            if spacing is None:
                line_px = size_px * 1.2
            elif hasattr(spacing, 'pt'):
                line_px = spacing.pt / 72 * PX
            else:
                line_px = size_px * spacing
            items.append(
                f'<div class="tx" style="left:{x}px;top:{y}px;width:{w}px;height:{h}px;'
                f'display:flex;flex-direction:column;justify-content:{va};'
                f'font-family:{family};font-size:{size_px:.2f}px;'
                f'line-height:{line_px:.2f}px;color:#{color};'
                f'font-weight:{"700" if bold else "400"};'
                f'font-style:{"italic" if italic else "normal"};'
                # Текст оборачиваем в свой блок: у голого текстового узла
                # внутри flex одиночная строка съезжает вниз, и превью
                # показывает её не там, где она окажется в PowerPoint.
                f'text-align:{align}"><div>{text}</div></div>'
            )

    parts.append(
        f'<div class="wrap" id="s{idx}"><div class="num">Слайд {idx}</div>'
        f'<div class="slide" style="{style}">' + ''.join(items) + '</div></div>'
    )

doc = f"""<!doctype html><html lang="ru"><head><meta charset="utf-8">
<title>Превью презентации</title><style>
body {{ margin:0; padding:24px; background:#3a4450; font-family:Calibri,Carlito,Arial,sans-serif; }}
.wrap {{ margin:0 auto 28px; width:{SW}px; }}
.num {{ color:#c9d3de; font-size:13px; margin-bottom:6px; }}
.slide {{ position:relative; width:{SW}px; height:{SH}px; overflow:hidden;
  box-shadow:0 8px 30px rgba(0,0,0,.35); }}
.shape, .tx, img {{ position:absolute; }}
.tx {{ white-space:pre-wrap; overflow-wrap:break-word; }}
img {{ object-fit:contain; }}
.over {{ outline:2px solid #ff4d4f; background:rgba(255,77,79,.12); }}
body.solo {{ padding:0; }}
body.solo .wrap {{ margin:0 auto; }}
body.solo .num {{ display:none; }}
</style></head><body class="{'solo' if ONLY else ''}">{''.join(parts)}
<script>
// Текст, который не влезает в свой блок, — самый частый дефект слайда.
// Меряем сами строки через Range: высота блока у крупных засечных
// шрифтов больше строки, и по ней получались ложные срабатывания.
window.overflows = [];
document.querySelectorAll('.slide').forEach((sl, i) => {{
  sl.querySelectorAll('.tx').forEach((el) => {{
    const range = document.createRange();
    range.selectNodeContents(el);
    const lines = [...range.getClientRects()].filter((r) => r.height > 0);
    if (!lines.length) return;
    const box = el.getBoundingClientRect();
    const top = Math.min(...lines.map((r) => r.top));
    const bottom = Math.max(...lines.map((r) => r.bottom));
    const wide = Math.max(...lines.map((r) => r.width));
    const high = bottom - top;
    if (high > box.height + 2 || wide > box.width + 2) {{
      el.classList.add('over');
      window.overflows.push(
        'слайд ' + (i + 1) + ': «' + el.textContent.slice(0, 42) + '…» ' +
        Math.round(high) + '×' + Math.round(wide) + ' в рамке ' +
        Math.round(box.height) + '×' + Math.round(box.width)
      );
    }}
  }});
}});

</script></body></html>"""

open(OUT, 'w', encoding='utf-8').write(doc)
print(f'слайдов: {len(prs.slides)}, размер: {SW:.0f}×{SH:.0f} px')
print('превью:', os.path.abspath(OUT))
if problems:
    print('\nвыходы за край:')
    for p in problems:
        print(' •', p)
else:
    print('за край слайда ничего не выходит')
